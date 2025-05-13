from docx import Document
from pptx import Presentation
import fitz
import os

def read_text(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        text = f.read()
    return text

    
def read_docx(file_path):
    doc = Document(file_path)
    text = '\n'.join([para.text for para in doc.paragraphs])
    return text

def read_pdf(file_path):
    text = ""
    pdf = fitz.open(file_path)
    for page in pdf:
        text += page.get_text()
    return text

def read_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


def extract_text(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    if ext == '.txt':
        return read_text(file_path)
    elif ext == '.docx':
        return read_docx(file_path)
    elif ext == '.pdf':
        return read_pdf(file_path)
    elif ext == '.pptx':
        return read_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")
